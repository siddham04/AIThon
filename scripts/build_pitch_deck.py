#!/usr/bin/env python3
"""
Build Helix hackathon pitch deck (.pptx) from PRESENTATION.md.

Usage (from repo root):
  pip install -r scripts/requirements-presentation.txt
  python scripts/build_pitch_deck.py

Output: docs/Helix-AI-Thon-Pitch.pptx (override with --out).
"""
from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path


def strip_markdown(text: str) -> str:
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\[(.+?)\]\([^)]+\)", r"\1", text)
    return text


def should_skip_body_line(line: str) -> bool:
    s = line.strip()
    if not s:
        return False
    if s.startswith("*(Insert") or s.startswith("*("):
        return True
    if "*(Insert screenshot" in s:
        return True
    return False


def parse_slides(md: str) -> list[tuple[str, str]]:
    slides: list[tuple[str, str]] = []
    for chunk in re.split(r"\n---\s*\n", md):
        chunk = chunk.strip()
        if not chunk:
            continue
        m = re.match(
            r"^###\s+Slide\s+\d+\s+[—-]\s*(.+?)\n+(.*)$",
            chunk,
            re.DOTALL | re.MULTILINE,
        )
        if not m:
            continue
        title = m.group(1).strip()
        body_raw = m.group(2).strip()
        lines = [ln for ln in body_raw.splitlines() if not should_skip_body_line(ln)]
        body = strip_markdown("\n".join(lines).strip())
        slides.append((title, body))
    return slides


def build_pptx(slides: list[tuple[str, str]], out_path: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError as e:
        print("Missing python-pptx. Run: pip install -r scripts/requirements-presentation.txt", file=sys.stderr)
        raise SystemExit(1) from e

    prs = Presentation()
    prs.slide_width = 9144000   # 10 in
    prs.slide_height = 6858000  # 7.5 in

    # Title slide
    layout_title = prs.slide_layouts[0]
    s0 = prs.slides.add_slide(layout_title)
    s0.shapes.title.text = "Helix"
    if s0.shapes.placeholders and len(s0.shapes.placeholders) > 1:
        sub = s0.placeholders[1]
        sub.text = "Intelligent SDLC Copilot — Requirement to delivery\nAI-Thon · Prototype Round 2"

    layout_content = prs.slide_layouts[1]
    for slide_title, body in slides:
        slide = prs.slides.add_slide(layout_content)
        slide.shapes.title.text = slide_title
        body_ph = slide.placeholders[1]
        tf = body_ph.text_frame
        paragraphs = [p.strip() for p in re.split(r"\n\s*\n", body) if p.strip()]
        if not paragraphs:
            paragraphs = [body] if body else ["(See repository documentation.)"]
        tf.text = paragraphs[0]
        p0 = tf.paragraphs[0]
        p0.level = 0
        p0.font.size = Pt(18 if len(paragraphs[0]) < 400 else 16)
        p0.space_after = Pt(10)
        for para in paragraphs[1:]:
            p = tf.add_paragraph()
            p.text = para
            p.level = 0
            p.font.size = Pt(18 if len(para) < 400 else 16)
            p.space_after = Pt(10)
        tf.word_wrap = True

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Wrote {out_path.resolve()} ({len(slides) + 1} slides including title).")


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    ap = argparse.ArgumentParser(description="Build .pptx from PRESENTATION.md")
    ap.add_argument(
        "--md",
        type=Path,
        default=root / "PRESENTATION.md",
        help="Source markdown",
    )
    ap.add_argument(
        "--out",
        type=Path,
        default=root / "docs" / "Helix-AI-Thon-Pitch.pptx",
        help="Output .pptx path",
    )
    args = ap.parse_args()
    md_text = args.md.read_text(encoding="utf-8")
    slides = parse_slides(md_text)
    if not slides:
        print("No slides parsed from markdown. Check PRESENTATION.md format.", file=sys.stderr)
        raise SystemExit(2)
    build_pptx(slides, args.out)


if __name__ == "__main__":
    main()
